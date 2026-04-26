"""
Build the bilingual ProjectEx walkthrough decks.

Generates two .pptx files in the parent ProjectEx directory:
  - ProjectEx-Walkthrough-TH.pptx
  - ProjectEx-Walkthrough-EN.pptx

Uses 16:9 widescreen layout, embeds the mockup screenshots from ../screenshots/.
"""

from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "screenshots"

# ---------- theme ----------
PURPLE = RGBColor(0x7B, 0x61, 0xFF)
DARK = RGBColor(0x18, 0x18, 0x1C)
LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
TEXT = RGBColor(0x22, 0x22, 0x28)
DIM = RGBColor(0x66, 0x66, 0x70)
GREEN = RGBColor(0x2D, 0xA8, 0x44)
BG_PANEL = RGBColor(0xEE, 0xEC, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, text, left, top, width, height, *, size=18, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_bullets(slide, items, left, top, width, height, *, size=16, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = TEXT
        run.font.name = font_name
        p.space_after = Pt(6)
    return box


def add_accent_bar(slide, top=Inches(0), left=Inches(0)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, SLIDE_W, Inches(0.18))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    return bar


def add_screenshot(slide, name, left, top, width):
    path = SHOTS / name
    pic = slide.shapes.add_picture(str(path), left, top, width=width)
    return pic


def add_footer(slide, text, page_num, total):
    add_text(slide, text, Inches(0.4), Inches(7.05), Inches(8), Inches(0.35),
             size=10, color=DIM)
    add_text(slide, f"{page_num} / {total}", Inches(12.4), Inches(7.05),
             Inches(0.6), Inches(0.35), size=10, color=DIM, align=PP_ALIGN.RIGHT)


def title_slide(prs, *, title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK)
    # purple block
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.4), Inches(2.4),
                                   Inches(0.18), Inches(2.5))
    block.line.fill.background()
    block.fill.solid()
    block.fill.fore_color.rgb = PURPLE
    add_text(slide, title, Inches(0.85), Inches(2.4), Inches(11.5), Inches(1.4),
             size=44, bold=True, color=LIGHT)
    add_text(slide, subtitle, Inches(0.85), Inches(3.9), Inches(11.5), Inches(1.0),
             size=22, color=RGBColor(0xC0, 0xC0, 0xCC))
    add_text(slide, "ProjectEx · Todo CLI",
             Inches(0.4), Inches(0.4), Inches(8), Inches(0.4),
             size=14, bold=True, color=PURPLE)
    add_text(slide, footer, Inches(0.4), Inches(6.9), Inches(12), Inches(0.4),
             size=12, color=RGBColor(0x99, 0x99, 0xA8))
    return slide


def content_slide(prs, *, title, subtitle=None, total_pages=12, page_num,
                  footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT)
    add_accent_bar(slide)
    add_text(slide, title, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             size=28, bold=True, color=TEXT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(1.05), Inches(12),
                 Inches(0.5), size=15, color=DIM)
    add_footer(slide, footer, page_num, total_pages)
    return slide


def step_slide(prs, *, step_num, title, image, bullets, page_num, total, footer,
               step_label="Step"):
    slide = content_slide(prs, title=f"{step_label} {step_num}: {title}",
                          page_num=page_num, total_pages=total, footer=footer)
    # screenshot on the right (~58% width)
    add_screenshot(slide, image, Inches(5.3), Inches(1.7), Inches(7.7))
    # bullets on the left
    add_bullets(slide, bullets, Inches(0.5), Inches(1.7), Inches(4.6),
                Inches(5.0), size=15)
    return slide


# ---------- TH content ----------
TH = {
    "filename": "ProjectEx-Walkthrough-TH.pptx",
    "footer": "CLAUDE-CODE-MANUAL · ProjectEx — Todo CLI Walkthrough",
    "slides": [
        {"kind": "title",
         "title": "สร้าง Todo CLI ด้วย Claude Code",
         "subtitle": "ตัวอย่างประกอบคู่มือ — สร้างจริง รันจริง คอมมิตจริง"},
        {"kind": "content", "title": "เป้าหมายและ Stack",
         "subtitle": "สิ่งที่คุณจะได้ลงมือทำในเวิร์กช็อปนี้",
         "body": [
             ("h", "เป้าหมาย"),
             ("b", "สร้าง Todo CLI ขนาดเล็ก: add / list / done / rm / edit"),
             ("b", "เน้น flow การทำงานกับ Claude Code มากกว่าตัวโค้ด"),
             ("b", "เห็น Plan Mode, Subagent, Hooks ทำงานจริง"),
             ("h", "Stack"),
             ("b", "Node.js (built-in only — ไม่มี dependencies)"),
             ("b", "Test runner: node:test (มากับ Node)"),
             ("b", "Storage: ไฟล์ JSON แบบเรียบ (gitignored)"),
             ("h", "ใช้เวลาประมาณ"),
             ("b", "45–60 นาที สำหรับมือใหม่"),
         ]},
        {"kind": "step", "step_num": 1, "title": "เริ่ม Claude Code",
         "image": "01-claude-start.png",
         "bullets": [
             "เข้าโฟลเดอร์โปรเจกต์ก่อน",
             "พิมพ์ `claude` เพื่อเปิด session",
             "Claude อ่าน `CLAUDE.md` อัตโนมัติ",
             "สังเกต cwd ที่หัวกล่อง — ต้องตรงกับโปรเจกต์",
             "Tip: Shift+Tab → Plan Mode สำหรับการเปลี่ยนหลายไฟล์",
         ]},
        {"kind": "step", "step_num": 2, "title": "เขียน CLAUDE.md ของโปรเจกต์",
         "image": "02-claude-md.png",
         "bullets": [
             "สั่งให้ Claude ดูโครงสร้างโค้ดแล้วร่าง CLAUDE.md",
             "ระบุ stack, conventions, สิ่งที่ห้ามทำ (don'ts)",
             "ไฟล์นี้จะถูกโหลดทุก session ในอนาคต",
             "ลด miscommunication เรื่อง style และข้อกำหนด",
             "ดูคู่มือบทที่ 7 — CLAUDE.md",
         ]},
        {"kind": "step", "step_num": 3, "title": "Plan Mode — วางแผนก่อนแก้",
         "image": "03-plan-mode.png",
         "bullets": [
             "สั่งงานที่กระทบหลายไฟล์ → ใช้ Plan Mode",
             "Claude ร่างแผนก่อน ไม่แตะไฟล์จนกว่าจะอนุมัติ",
             "ตรวจสอบขอบเขต: ตรงตามที่อยากได้ไหม?",
             "อนุมัติด้วย y → เริ่ม implement",
             "ดูคู่มือบทที่ 14 — Context Management",
         ]},
        {"kind": "step", "step_num": 4, "title": "ลงมือเขียนโค้ด — ดู Diff",
         "image": "04-edit-diff.png",
         "bullets": [
             "Claude แก้ทีละไฟล์ตามแผน",
             "แสดง diff สีเขียว (+) / แดง (–) ก่อนเขียนจริง",
             "Hook PostToolUse จะรัน `npm test` หลัง Edit แต่ละครั้ง",
             "ถ้า test แดง → fix ก่อนไปต่อ",
             "ดูคู่มือบทที่ 10 — Hooks",
         ]},
        {"kind": "step", "step_num": 5, "title": "รัน Test ทั้งหมด",
         "image": "05-tests.png",
         "bullets": [
             "เรียก `/test-all` — slash command ที่เราสร้างเอง",
             "ใช้ allowed-tools จำกัดให้รันได้แค่ `npm test`",
             "Claude สรุปผล PASS/FAIL ในบรรทัดเดียว",
             "ถ้าแดง: Claude ระบุสาเหตุและเสนอ fix (ไม่ apply เอง)",
             "ดูคู่มือบทที่ 11 — Skills / Slash Commands",
         ]},
        {"kind": "step", "step_num": 6, "title": "ใช้ Subagent ตรวจโค้ด",
         "image": "06-review.png",
         "bullets": [
             "เรียก subagent `reviewer` ที่ config ไว้",
             "ทำงานในบริบทแยก — ไม่เปลือง context หลัก",
             "tools จำกัด: Read, Grep, git diff/log เท่านั้น",
             "รายงาน LGTM / NEEDS WORK พร้อม punch list",
             "ดูคู่มือบทที่ 12 — Subagents",
         ]},
        {"kind": "step", "step_num": 7, "title": "Commit ผ่าน Claude",
         "image": "07-commit.png",
         "bullets": [
             "Claude ร่าง commit message จาก diff",
             "เน้น \"why\" มากกว่า \"what\"",
             "ใช้ Co-Authored-By: Claude…",
             "Permission allow `Bash(git commit:*)` ไว้แล้ว",
             "ดูคู่มือบทที่ 15 — Git Integration",
         ]},
        {"kind": "content", "title": "Hooks & Skills ที่ใช้ในตัวอย่างนี้",
         "body": [
             ("h", "Hook: PostToolUse (Edit|Write)"),
             ("b", "รัน `npm test` อัตโนมัติหลังแก้ไฟล์"),
             ("b", "ตั้งใน `.claude/settings.json`"),
             ("h", "Slash Command: /test-all"),
             ("b", "ไฟล์: `.claude/commands/test-all.md`"),
             ("b", "allowed-tools จำกัดให้รันได้เฉพาะ `npm test`"),
             ("h", "Subagent: reviewer"),
             ("b", "ไฟล์: `.claude/agents/reviewer.md`"),
             ("b", "model: sonnet, tools: Read/Grep/git diff/git log"),
             ("h", "Permissions"),
             ("b", "allow: npm test, node *, git status/diff/log/add/commit"),
             ("b", "deny: rm -rf, git push --force, write todos.json"),
         ]},
        {"kind": "content", "title": "บทเรียนที่ได้",
         "body": [
             ("h", "ทำอะไรได้ดี"),
             ("b", "Plan Mode ช่วยลด rework เวลาแก้หลายไฟล์"),
             ("b", "Subagent แยก context — main session ไม่รก"),
             ("b", "Hooks ทำให้ test ไหลเป็นส่วนหนึ่งของ workflow"),
             ("h", "ข้อควรระวัง"),
             ("b", "อย่า allow `Bash(*)` กว้างเกินไป"),
             ("b", "Hook ที่ช้าจะทำให้ทุกการแก้ไฟล์ช้าตาม"),
             ("b", "CLAUDE.md ที่ยาวเกินจะกิน context — เก็บแค่จำเป็น"),
             ("h", "ขั้นถัดไป"),
             ("b", "ลองเพิ่ม `todo prio <id> <high|low>` ด้วย Plan Mode"),
             ("b", "ตั้ง Scheduled Task ทบทวน todos รายวัน"),
         ]},
        {"kind": "content", "title": "อ่านต่อ — ลิงก์ไปยังคู่มือ",
         "body": [
             ("h", "บทที่เกี่ยวข้องโดยตรง"),
             ("b", "บท 7 — CLAUDE.md (คำสั่งถาวรของโปรเจกต์)"),
             ("b", "บท 10 — Hooks (Event Handler)"),
             ("b", "บท 11 — Skills / Slash Commands"),
             ("b", "บท 12 — Subagents"),
             ("b", "บท 14 — Context Management (Plan Mode)"),
             ("b", "บท 15 — Git Integration"),
             ("h", "อ่านเพิ่ม"),
             ("b", "บท 28 — Tutorial Day 2: สร้าง Todo App ใน 1 ชั่วโมง"),
             ("b", "บท 30 — Cookbook (40+ Recipes)"),
             ("b", "บท 32 — Security & Privacy Best Practices"),
             ("h", "ไฟล์โปรเจกต์ตัวอย่างนี้"),
             ("b", "ProjectEx/todo-app — โค้ดรันได้ + tests"),
             ("b", "ProjectEx/claude-config — CLAUDE.md, hooks, agents"),
         ]},
    ],
    "step_label": "ขั้นที่",
}

# ---------- EN content ----------
EN = {
    "filename": "ProjectEx-Walkthrough-EN.pptx",
    "footer": "CLAUDE-CODE-MANUAL · ProjectEx — Todo CLI Walkthrough",
    "slides": [
        {"kind": "title",
         "title": "Build a Todo CLI with Claude Code",
         "subtitle": "Manual companion — built, run, and committed for real"},
        {"kind": "content", "title": "Goal & Stack",
         "subtitle": "What you'll build in this walkthrough",
         "body": [
             ("h", "Goal"),
             ("b", "A small Todo CLI: add / list / done / rm / edit"),
             ("b", "Focus on the Claude Code workflow, not the code"),
             ("b", "See Plan Mode, Subagents, and Hooks in action"),
             ("h", "Stack"),
             ("b", "Node.js — built-ins only, zero dependencies"),
             ("b", "Test runner: node:test (ships with Node)"),
             ("b", "Storage: flat JSON file (gitignored)"),
             ("h", "Time"),
             ("b", "About 45–60 minutes for a beginner"),
         ]},
        {"kind": "step", "step_num": 1, "title": "Start a Claude session",
         "image": "01-claude-start.png",
         "bullets": [
             "`cd` into your project directory first",
             "Run `claude` to start the session",
             "Claude auto-loads `CLAUDE.md`",
             "Verify cwd in the welcome banner",
             "Tip: Shift+Tab → Plan Mode for multi-file changes",
         ]},
        {"kind": "step", "step_num": 2, "title": "Author CLAUDE.md",
         "image": "02-claude-md.png",
         "bullets": [
             "Ask Claude to scan the repo and draft CLAUDE.md",
             "Capture stack, conventions, and don'ts",
             "This file loads automatically every session",
             "Reduces miscommunication on style and constraints",
             "See chapter 7 — CLAUDE.md",
         ]},
        {"kind": "step", "step_num": 3, "title": "Plan Mode — plan before edits",
         "image": "03-plan-mode.png",
         "bullets": [
             "For multi-file changes, switch to Plan Mode",
             "Claude drafts a plan and waits for approval",
             "Review scope before any file is touched",
             "Approve with `y` to start implementing",
             "See chapter 14 — Context Management",
         ]},
        {"kind": "step", "step_num": 4, "title": "Implement & inspect the diff",
         "image": "04-edit-diff.png",
         "bullets": [
             "Claude edits one file at a time, per the plan",
             "Diff shown with + (green) / – (red) before write",
             "PostToolUse hook runs `npm test` after each edit",
             "If tests fail — fix before continuing",
             "See chapter 10 — Hooks",
         ]},
        {"kind": "step", "step_num": 5, "title": "Run all tests",
         "image": "05-tests.png",
         "bullets": [
             "Invoke our custom slash command `/test-all`",
             "`allowed-tools` restricts it to `npm test` only",
             "Claude reports PASS/FAIL in one line",
             "On failure: it suggests a fix without applying it",
             "See chapter 11 — Skills / Slash Commands",
         ]},
        {"kind": "step", "step_num": 6, "title": "Subagent code review",
         "image": "06-review.png",
         "bullets": [
             "Dispatch the configured `reviewer` subagent",
             "Runs in an isolated context — main session stays clean",
             "Tools: Read, Grep, git diff/log only",
             "Returns LGTM / NEEDS WORK with a punch list",
             "See chapter 12 — Subagents",
         ]},
        {"kind": "step", "step_num": 7, "title": "Commit through Claude",
         "image": "07-commit.png",
         "bullets": [
             "Claude drafts the commit message from the diff",
             "Focus on \"why\" over \"what\"",
             "Adds Co-Authored-By: Claude…",
             "`Bash(git commit:*)` is pre-allowed in settings",
             "See chapter 15 — Git Integration",
         ]},
        {"kind": "content", "title": "Hooks & Skills used here",
         "body": [
             ("h", "Hook: PostToolUse (Edit|Write)"),
             ("b", "Auto-runs `npm test` after every edit"),
             ("b", "Configured in `.claude/settings.json`"),
             ("h", "Slash command: /test-all"),
             ("b", "File: `.claude/commands/test-all.md`"),
             ("b", "`allowed-tools` restricts it to `npm test`"),
             ("h", "Subagent: reviewer"),
             ("b", "File: `.claude/agents/reviewer.md`"),
             ("b", "model: sonnet · tools: Read/Grep/git diff/git log"),
             ("h", "Permissions"),
             ("b", "allow: npm test, node *, git status/diff/log/add/commit"),
             ("b", "deny: rm -rf, git push --force, writes to todos.json"),
         ]},
        {"kind": "content", "title": "Lessons learned",
         "body": [
             ("h", "What worked well"),
             ("b", "Plan Mode reduced rework on multi-file changes"),
             ("b", "Subagent isolation kept the main context clean"),
             ("b", "Hooks made tests part of the workflow, not an afterthought"),
             ("h", "Watch out for"),
             ("b", "Don't allow `Bash(*)` — keep the allowlist narrow"),
             ("b", "Slow hooks slow every single edit"),
             ("b", "Long CLAUDE.md eats context — keep it lean"),
             ("h", "Next steps"),
             ("b", "Add `todo prio <id> <high|low>` via Plan Mode"),
             ("b", "Schedule a daily todo-review task"),
         ]},
        {"kind": "content", "title": "Further reading — manual links",
         "body": [
             ("h", "Directly relevant chapters"),
             ("b", "Ch 7 — CLAUDE.md (persistent project instructions)"),
             ("b", "Ch 10 — Hooks (event handlers)"),
             ("b", "Ch 11 — Skills / Slash Commands"),
             ("b", "Ch 12 — Subagents"),
             ("b", "Ch 14 — Context Management (Plan Mode)"),
             ("b", "Ch 15 — Git Integration"),
             ("h", "Going deeper"),
             ("b", "Ch 28 — Tutorial Day 2: Build a Todo App in 1 hour"),
             ("b", "Ch 30 — Cookbook (40+ recipes)"),
             ("b", "Ch 32 — Security & Privacy Best Practices"),
             ("h", "This example's source"),
             ("b", "ProjectEx/todo-app — runnable code + tests"),
             ("b", "ProjectEx/claude-config — CLAUDE.md, hooks, agents"),
         ]},
    ],
    "step_label": "Step",
}


def render_body(slide, body, *, left=Inches(0.5), top=Inches(1.4),
                width=Inches(12.3)):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for kind, text in body:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if kind == "h":
            p.space_before = Pt(8)
            p.space_after = Pt(2)
            run = p.add_run()
            run.text = text
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = PURPLE
            run.font.name = "Calibri"
        else:  # bullet
            p.space_after = Pt(2)
            run = p.add_run()
            run.text = "  •  " + text
            run.font.size = Pt(15)
            run.font.color.rgb = TEXT
            run.font.name = "Calibri"


def build_deck(spec):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(spec["slides"])
    footer = spec["footer"]
    step_label = spec["step_label"]

    for i, sl in enumerate(spec["slides"], start=1):
        kind = sl["kind"]
        if kind == "title":
            title_slide(prs, title=sl["title"], subtitle=sl["subtitle"],
                        footer=footer)
        elif kind == "step":
            step_slide(prs, step_num=sl["step_num"], title=sl["title"],
                       image=sl["image"], bullets=sl["bullets"],
                       page_num=i, total=total, footer=footer,
                       step_label=step_label)
        else:
            slide = content_slide(prs, title=sl["title"],
                                  subtitle=sl.get("subtitle"),
                                  page_num=i, total_pages=total, footer=footer)
            render_body(slide, sl["body"])

    out = ROOT / spec["filename"]
    prs.save(out)
    print(f"  wrote {out}")


def main():
    print("Building decks…")
    build_deck(TH)
    build_deck(EN)
    print("Done.")


if __name__ == "__main__":
    main()
